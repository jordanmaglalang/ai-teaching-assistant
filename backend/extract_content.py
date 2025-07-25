from docx import Document
import fitz 
from pptx import Presentation

def get_doc_content(content):
    doc = Document(content)
    text_content = ""
    for paragraph in doc.paragraphs:
        text_content += paragraph.text + "\n"  # add each paragraph text with a newline
        print(text_content)
    return text_content
    
def get_text_from_pdf(path):
    pdf_bytes = path.read()
    doc = fitz.open(stream = pdf_bytes, filetype = "pdf")  # open the PDF file
    text_content = ""
    for page in doc:
        text_content += page.get_text() + "\n"  # extract text from each page
    print(text_content)
    return text_content
def get_text_from_ppt(path):
    prs = Presentation(path)
    text_content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    #print(text_content)
    return text_content
#doc_path = r"C:\Users\jorda\OneDrive\Documents\testing content ta-ai.docx"
#pdf_path = r"C:\Users\jorda\Downloads\cmsc313-assembly-nasm-intro.pdf"
pptx_path = r"/Users/jordanmaglalang/Library/CloudStorage/OneDrive-Personal/Documents/01-Intro.pptx"
pdf_path =r"/Users/jordanmaglalang/Downloads/Stack.pdf"
#get_doc_content(doc_path)

